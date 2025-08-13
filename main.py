#install required libraries
# pip install google-generativeai python-pptx Pillow  

import google.generativeai as genai
import sys
import os
from pptx import Presentation
from PIL import Image
import io

API_KEY = 'YOUR-GEMINI-API-KEY'

try:
    genai.configure(api_key=API_KEY)
except Exception as e:
    print(f"Error configuring Gemini API: {e}")
    print("Please ensure you have set a valid API_KEY.")
    sys.exit(1)

def analyze_presentation(pptx_path):
    """
    Analyzes a PowerPoint presentation for inconsistencies using the Gemini API.

    Args:
        pptx_path (str): The file path to the .pptx presentation.
    """
    if not os.path.exists(pptx_path):
        print(f"Error: File not found at '{pptx_path}'")
        return

    print(f"🔍 Analyzing '{os.path.basename(pptx_path)}'...")
    
    try:
        prs = Presentation(pptx_path)
        model = genai.GenerativeModel('gemini-2.5-flash')
        all_slides_content = []

        # 1. Extract content from all slides first
        for i, slide in enumerate(prs.slides):
            slide_number = i + 1
            slide_content = {"text": f"--- Slide {slide_number} Text --- \n"}
            
            # Extract text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_content["text"] += shape.text + "\n"
            
            # Extract images
            slide_images = []
            for shape in slide.shapes:
                if hasattr(shape, 'image'):
                    image_bytes = shape.image.blob
                    try:
                        img = Image.open(io.BytesIO(image_bytes))
                        slide_images.append(img)
                    except Exception as e:
                        print(f"Warning: Could not process an image on slide {slide_number}: {e}")
            
            slide_content["images"] = slide_images
            all_slides_content.append(slide_content)
        
        # 2. Build the full prompt for the model
        full_prompt_parts = [
            """
            You are an expert business analyst. Your task is to analyze the content of the following presentation slides. 
            Review all the provided text and images from the entire presentation and identify any discrepancies or mismatches. 
            Focus on these types of inconsistencies:
            1.  **Conflicting Numerical Data**: Revenue figures, statistics, or percentages that don't add up or contradict each other across different slides.
            2.  **Contradictory Textual Claims**: Statements that conflict with each other (e.g., describing a market as "highly competitive" on one slide and having "few competitors" on another).
            3.  **Timeline Mismatches**: Dates, project phases, or forecasts that are logically inconsistent or conflict.
            
            Please provide a structured summary of your findings. For each discrepancy found, reference the slide numbers involved and clearly explain the nature of the issue.
            If no inconsistencies are found, simply state "No inconsistencies found."

            Here is the content from all the slides:
            """
        ]

        # Add text and images from each slide to the prompt
        for i, content in enumerate(all_slides_content):
            slide_number = i + 1
            full_prompt_parts.append(f"\n--- Slide {slide_number} Content ---")
            full_prompt_parts.append(content["text"])
            if content["images"]:
                full_prompt_parts.append(f"Images from Slide {slide_number}:")
                for img in content["images"]:
                    full_prompt_parts.append(img)
        
        # 3. Call the Gemini API once with all content
        print("\nSending all slides to Gemini for a holistic analysis... 🧠")
        response = model.generate_content(full_prompt_parts, stream=False)
        
        # 4. Print the final report
        print("\n--- 📊 Final Analysis Report ---")
        print(response.text)
        print("---------------------------------")
        print("\n✅ Analysis complete.")

    except Exception as e:
        print(f"\nAn error occurred during analysis: {e}")
        print("This could be due to a corrupted file, an invalid API key, or a network issue.")


# --- Script Execution ---

if __name__ == "__main__":
    if len(sys.argv) != 2:
        print("Usage: python analyze_ppt.py <path_to_presentation.pptx>")
        sys.exit(1)
    
    file_path = sys.argv[1]
    analyze_presentation(file_path)

